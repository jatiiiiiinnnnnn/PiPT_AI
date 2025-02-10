import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List, Dict
from dataclasses import dataclass
import json
import logging
import os
from dotenv import load_dotenv

st.set_page_config(
    page_title="AI PowerPoint Generator",
    page_icon="https://img.icons8.com/?size=100&id=52564&format=png&color=000000",
    layout="wide"
)

# First, configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Get API key from Streamlit secrets and configure Gemini
# Configure Gemini directly with the API key from secrets
try:
    genai.configure(api_key=st.secrets["GOOGLE_API_KEY"])  # Direct use of secret, no variable
except Exception as e:
    st.error(f"Failed to initialize Gemini API: {str(e)}")
    logger.error(f"API configuration error: {str(e)}")
    st.stop()

# Debug button to verify secret is loaded
if st.button("Debug: Check API Key"):
    st.write("First few characters of API key:", st.secrets["GOOGLE_API_KEY"][:5] + "...")



#load_dotenv()
#API_KEY = os.getenv("GOOGLE_API_KEY")
#if not API_KEY:
#    raise ValueError("API key not found. Make sure .env contains GOOGLE_API_KEY and is properly loaded.")

## Configure Google Gemini API with the key
#genai.configure(api_key=API_KEY)



@dataclass
class StyleSettings:
    font_style: str
    font_size: str
    color_theme: str


# Constants
FONT_STYLES = ["Arial", "Calibri", "Times New Roman", "Helvetica", "Georgia"]
FONT_SIZES = {
    "Small": {"title": 32, "subtitle": 24, "body": 18},
    "Medium": {"title": 40, "subtitle": 28, "body": 20},
    "Large": {"title": 44, "subtitle": 32, "body": 24}
}
COLOR_THEMES = {
    "Professional Blue": {"primary": RGBColor(0, 75, 150), "secondary": RGBColor(255, 255, 255)},
    "Forest Green": {"primary": RGBColor(34, 139, 34), "secondary": RGBColor(255, 255, 255)},
    "Classic Gray": {"primary": RGBColor(80, 80, 80), "secondary": RGBColor(255, 255, 255)},
    "Deep Purple": {"primary": RGBColor(75, 0, 130), "secondary": RGBColor(255, 255, 255)},
    "Elegant Black": {"primary": RGBColor(0, 0, 0), "secondary": RGBColor(255, 255, 255)}
}

SLIDE_LAYOUTS = {
    "Title Slide": 0,
    "Content": 1,
    "Section Header": 2,
    "Two Content": 3,
    "Comparison": 4
}


@st.cache_resource
def initialize_genai():
    """Initialize the Gemini model with caching"""
     """Initialize the Gemini model with caching"""
    try:
        api_key = st.secrets["GOOGLE_API_KEY"]  # Fetch API key from Streamlit secrets
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-pro')
        return model
    except Exception as e:
        logger.error(f"Failed to initialize Gemini model: {e}")
        raise


def generate_content(prompt: str, model) -> str:
    """Generate content using the Gemini model"""
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        logger.error(f"Error generating content: {e}")
        raise


def generate_presentation_content(topic: str, num_slides: int, purpose: str, audience: str, model) -> List[Dict]:
    """Generate more detailed presentation content using the Gemini model"""
    try:
        prompt = f"""Create a compelling and engaging presentation outline for '{topic}' with exactly {num_slides - 1} content slides.
        Purpose: {purpose}
        Target Audience: {audience}

        Follow these guidelines for high-impact content:
        1. Make the title slide captivating and memorable
        2. Each content slide should:
           - Have a clear, action-oriented title
           - Include a thought-provoking subtitle that reinforces the main message
           - Present 3-4 key points with specific examples and data
           - Tell a coherent story that flows from the previous slide
        3. Include real-world applications or case studies
        4. End with a strong call-to-action or key takeaways

        Structure the content to maintain audience engagement through:
        - Opening with a hook or surprising fact
        - Building narrative tension
        - Using concrete examples and statistics
        - Including interactive elements or discussion points
        - Connecting points to audience's needs and interests

        Format as a JSON array with:
        {{
            "title": "Compelling Title",
            "subtitle": "Engaging Subtitle or Key Message",
            "content": [
                "Detailed point with specific example or data",
                "Real-world application or case study",
                "Actionable insight or key learning"
            ],
            "examples": ["Specific example", "Case study", "Statistic"],
            "image_description": "Detailed description for relevant, impactful image"
        }}"""

        response = generate_content(prompt, model)
        cleaned_response = clean_json_response(response)
        return json.loads(cleaned_response)

    except Exception as e:
        logger.error(f"Error in detailed content generation: {e}")
        raise


def clean_json_response(response: str) -> str:
    """Clean the API response to ensure valid JSON"""
    try:
        # Remove any potential markdown formatting
        response = response.replace("```json", "").replace("```", "")

        # Remove any leading/trailing whitespace
        response = response.strip()

        # If response starts with additional text, try to find the start of the JSON array
        if not response.startswith("["):
            start_idx = response.find("[")
            if start_idx != -1:
                response = response[start_idx:]

        # If response has additional text after the JSON array, try to find the end
        if not response.endswith("]"):
            end_idx = response.rfind("]")
            if end_idx != -1:
                response = response[:end_idx + 1]

        # Validate JSON structure
        json.loads(response)  # This will raise JSONDecodeError if invalid
        return response

    except json.JSONDecodeError as e:
        logger.error(f"Failed to clean JSON response: {e}")
        # If cleaning fails, try to construct a basic valid JSON response
        fallback_content = [
            {
                "title": "Technical Difficulties",
                "content": [
                    "We apologize for the technical issue.",
                    "Please try regenerating the presentation.",
                    "If the problem persists, try simplifying your topic."
                ]
            }
        ]
        return json.dumps(fallback_content)

def apply_slide_styling(slide, style_settings: StyleSettings):
    """Apply styling to a slide"""
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            text_frame = shape.text_frame
            text_frame.word_wrap = True

            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = style_settings.font_style

                    if hasattr(shape, "name") and "title" in shape.name.lower():
                        run.font.size = Pt(FONT_SIZES[style_settings.font_size]["title"])
                    elif hasattr(shape, "name") and "subtitle" in shape.name.lower():
                        run.font.size = Pt(FONT_SIZES[style_settings.font_size]["subtitle"])
                    else:
                        run.font.size = Pt(FONT_SIZES[style_settings.font_size]["body"])

                    run.font.color.rgb = COLOR_THEMES[style_settings.color_theme]["primary"]


def create_powerpoint(slides_content: List[Dict], style_settings: StyleSettings) -> str:
    """Create PowerPoint presentation"""
    prs = Presentation()

    # Create title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUTS["Title Slide"]])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = slides_content[0]['title']
    subtitle.text = slides_content[0]['content'][0] if slides_content[0]['content'] else ""

    apply_slide_styling(title_slide, style_settings)

    # Create content slides
    for slide_content in slides_content[1:]:
        content_slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUTS["Content"]])

        title = content_slide.shapes.title
        title.text = slide_content['title']

        body = content_slide.placeholders[1]
        tf = body.text_frame
        tf.clear()  # Clear existing text

        for point in slide_content['content']:
            p = tf.add_paragraph()
            p.text = str(point)  # Convert to string in case of non-string content
            p.level = 0

        apply_slide_styling(content_slide, style_settings)

    output_path = "generated_presentation.pptx"
    prs.save(output_path)
    return output_path


def main():
   
    st.image("https://img.icons8.com/?size=100&id=52564&format=png&color=000000",width=50)
    st.title("AI PowerPoint Generator")
    st.write("Generate professional presentations using AI")

    try:
        model = initialize_genai()
    except Exception as e:
        st.error(f"Failed to initialize the AI model. Please check your API key configuration. Error: {str(e)}")
        return

    # Sidebar styling options
    st.sidebar.title("Presentation Styling")
    style_settings = StyleSettings(
        font_style=st.sidebar.selectbox("Font Style", FONT_STYLES),
        font_size=st.sidebar.selectbox("Font Size", list(FONT_SIZES.keys())),
        color_theme=st.sidebar.selectbox("Color Theme", list(COLOR_THEMES.keys()))
    )

    # Main content area
    col1, col2 = st.columns([2, 1])

    with col1:
        topic = st.text_input("Presentation Topic:", placeholder="e.g., Introduction to AI")

        col_left, col_right = st.columns(2)
        with col_left:
            num_slides = st.slider("Number of slides:", min_value=3, max_value=10, value=5)
            purpose = st.selectbox("Purpose:", ["Informative", "Persuasive", "Educational", "Business"])

        with col_right:
            audience = st.text_input("Target Audience:", placeholder="e.g., Business professionals")

        if st.button("Generate Presentation", type="primary"):
            if topic and audience:
                try:
                    with st.spinner("Generating your presentation..."):
                        slides_content = generate_presentation_content(
                            topic, num_slides, purpose, audience, model
                        )

                        output_path = create_powerpoint(slides_content, style_settings)

                        with open(output_path, "rb") as file:
                            st.download_button(
                                label="📥 Download Presentation",
                                data=file,
                                file_name=f"{topic.lower().replace(' ', '_')}_presentation.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )

                        # Preview content
                        with col2:
                            st.success("✨ Presentation generated successfully!")
                            st.write("### Content Preview:")
                            for i, slide in enumerate(slides_content):
                                with st.expander(f"Slide {i + 1}: {slide['title']}"):
                                    for point in slide['content']:
                                        st.write(f"• {point}")

                except Exception as e:
                    st.error(f"An error occurred: {str(e)}")
            else:
                st.warning("Please fill in all required fields.")


if __name__ == "__main__":
    main()
